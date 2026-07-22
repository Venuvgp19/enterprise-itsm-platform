import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette (Dark Modern ITSM Theme)
    BG_DARK = RGBColor(15, 23, 42)       # Slate 900
    CARD_BG = RGBColor(30, 41, 59)      # Slate 800
    TEXT_WHITE = RGBColor(248, 250, 252) # Slate 50
    TEXT_MUTED = RGBColor(148, 163, 184)# Slate 400
    ACCENT_INDIGO = RGBColor(129, 140, 248) # Indigo 400
    ACCENT_PURPLE = RGBColor(192, 132, 252) # Purple 400
    ACCENT_EMERALD = RGBColor(52, 211, 153) # Emerald 400
    ACCENT_AMBER = RGBColor(251, 191, 36)  # Amber 400
    BORDER_COLOR = RGBColor(51, 65, 85)   # Slate 700

    blank_layout = prs.slide_layouts[6]

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="ENTERPRISE AGENTIC AI ITSM"):
        # Category Badge
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = ACCENT_INDIGO

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
        return shape

    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide1 = prs.slides.add_slide(blank_layout)
    add_bg(slide1)

    # Decorative Gradient Card Background
    add_card(slide1, 1.0, 1.2, 11.333, 5.1, bg_color=RGBColor(24, 32, 52), border_color=ACCENT_INDIGO)

    tb = slide1.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Enterprise Agentic AI ITSM Platform"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p2 = tf.add_paragraph()
    p2.text = "100% Autonomous Ticket Routing & Knowledge Base SOP Synthesis"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_PURPLE

    tb_desc = slide1.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(10.333), Inches(2.0))
    tf_desc = tb_desc.text_frame
    tf_desc.word_wrap = True
    
    bullets = [
      "🤖 Autonomous Ticket Routing: NVIDIA Nemotron 3 Ultra 550B LLM Engine",
      "🧠 Agentic Knowledge Base Generator: Meta Llama 3.3 70B Instruct Engine",
      "🔌 Model Context Protocol (MCP): Native stdio transport with Auto JWT Auth Flow",
      "💾 Resilient Persistence Layer: 1,000 Incidents & SOPs saved to disk zero data loss"
    ]
    for b in bullets:
        pb = tf_desc.add_paragraph()
        pb.text = b
        pb.font.size = Pt(14)
        pb.font.color.rgb = TEXT_MUTED

    # ==================== SLIDE 2: EXECUTIVE SUMMARY ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_bg(slide2)
    add_header(slide2, "Executive Summary & Core Platform Capabilities")

    cards_s2 = [
        ("Zero-Human Ticket Triage", "100% LLM-driven incident classification & technician routing using NVIDIA Nemotron 3 550B with backoff protection.", ACCENT_INDIGO),
        ("Automated SOP Synthesis", "Scans diagnostic work notes across 1,000 incidents in 10-ticket chunks using Meta Llama 3.3 70B Instruct.", ACCENT_PURPLE),
        ("MCP Server Integration", "Stdio transport protocol with automatic JWT auth flow, empowering external AI agents to manage ITSM entities.", ACCENT_EMERALD),
        ("Strict Deduplication", "Semantic token overlap check (>50%) ensures 100% unique problem solutions with persistent analyzed incident tracking.", ACCENT_AMBER),
    ]

    lefts = [0.8, 6.8, 0.8, 6.8]
    tops = [1.6, 1.6, 4.3, 4.3]

    for idx, (title, desc, accent) in enumerate(cards_s2):
        add_card(slide2, lefts[idx], tops[idx], 5.7, 2.4)
        tb_card = slide2.shapes.add_textbox(Inches(lefts[idx] + 0.3), Inches(tops[idx] + 0.3), Inches(5.1), Inches(1.8))
        tf_card = tb_card.text_frame
        tf_card.word_wrap = True
        
        p_ct = tf_card.paragraphs[0]
        p_ct.text = title
        p_ct.font.size = Pt(18)
        p_ct.font.bold = True
        p_ct.font.color.rgb = accent

        p_cd = tf_card.add_paragraph()
        p_cd.text = desc
        p_cd.font.size = Pt(13)
        p_cd.font.color.rgb = TEXT_MUTED

    # ==================== SLIDE 3: ARCHITECTURE DIAGRAM ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_bg(slide3)
    add_header(slide3, "End-to-End Multi-Agent System Architecture")

    # Layer Cards
    layers = [
        ("CLIENT LAYER", "Next.js Frontend (Port 3000)  |  MCP Stdio Client (itsm-mcp-server)", 1.6, ACCENT_INDIGO),
        ("ORCHESTRATION LAYER", "🤖 Continuous AI Ticket Router  |  🧠 Agentic Knowledge Synthesizer", 2.9, ACCENT_PURPLE),
        ("BACKEND API LAYER", "NestJS REST API (Port 4000)  |  AuthModule  |  IncidentModule  |  KnowledgeModule", 4.2, ACCENT_EMERALD),
        ("AI MODEL & DISK LAYER", "NVIDIA NIM API (Nemotron 3 550B & Llama 3.3 70B)  |  incidents.json & analyzed_incidents.json", 5.5, ACCENT_AMBER),
    ]

    for title, desc, top_y, accent in layers:
        add_card(slide3, 0.8, top_y, 11.733, 1.1)
        tb_l = slide3.shapes.add_textbox(Inches(1.1), Inches(top_y + 0.15), Inches(11.1), Inches(0.8))
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        
        p_lt = tf_l.paragraphs[0]
        p_lt.text = title
        p_lt.font.size = Pt(14)
        p_lt.font.bold = True
        p_lt.font.color.rgb = accent

        p_ld = tf_l.add_paragraph()
        p_ld.text = desc
        p_ld.font.size = Pt(12)
        p_ld.font.color.rgb = TEXT_WHITE

    # ==================== SLIDE 4: TICKET ROUTER ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_bg(slide4)
    add_header(slide4, "Autonomous AI Ticket Router (NVIDIA Nemotron 3 550B)")

    add_card(slide4, 0.8, 1.6, 5.7, 5.1)
    tb_r1 = slide4.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.7))
    tf_r1 = tb_r1.text_frame
    tf_r1.word_wrap = True
    
    p = tf_r1.paragraphs[0]
    p.text = "Router Mechanics & Resilience"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO

    r_bullets = [
        "• Unassigned Queue Worker: Scans DB every 10,000ms for unassigned incident tickets.",
        "• 100% LLM-Driven: Zero fallback to rule engine, ensuring true agentic reasoning.",
        "• Sequential Pacing: Processes tickets one-by-one with 2.0s pauses to prevent rate limits.",
        "• Backoff Retries: Exponential retry backoff (3s, 6s, 9s, 12s) handling HTTP 429/503 limits.",
        "• Audit Logging: Extracts extended thinking traces into work notes and ai-router.log."
    ]
    for b in r_bullets:
        pb = tf_r1.add_paragraph()
        pb.text = b
        pb.font.size = Pt(12)
        pb.font.color.rgb = TEXT_MUTED

    add_card(slide4, 6.8, 1.6, 5.7, 5.1)
    tb_r2 = slide4.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.7))
    tf_r2 = tb_r2.text_frame
    tf_r2.word_wrap = True
    
    p = tf_r2.paragraphs[0]
    p.text = "Sample LLM Routing Payload"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD

    json_sample = """{
  "ticketId": "INC0000008",
  "targetGroup": "Network Ops",
  "assignedTechnician": "Alex Rivera (Network Lead)",
  "confidenceScore": 95,
  "routedBy": "NVIDIA_NEMOTRON_LLM",
  "reasoningText": "Analyzed core router latency log. Matched Network Ops team lead for BGP flush.",
  "recommendedResolutionCode": "Network - BGP & Interface Reset"
}"""
    pb = tf_r2.add_paragraph()
    pb.text = json_sample
    pb.font.size = Pt(11)
    pb.font.bold = True
    pb.font.color.rgb = TEXT_WHITE

    # ==================== SLIDE 5: KNOWLEDGE BASE SYNTHESIZER ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_bg(slide5)
    add_header(slide5, "Agentic Knowledge Base Synthesizer (Meta Llama 3.3 70B)")

    add_card(slide5, 0.8, 1.6, 5.7, 5.1)
    tb_k1 = slide5.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.7))
    tf_k1 = tb_k1.text_frame
    tf_k1.word_wrap = True
    
    p = tf_k1.paragraphs[0]
    p.text = "Chunked Processing & Deduplication"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE

    k_bullets = [
        "• 10-Incident Chunking: Processes unanalyzed tickets in 10-ticket batches for optimal LLM context window use.",
        "• Persistent Incident Tracking: Marks analyzed IDs in analyzed_incidents.json so no incident is ever re-analyzed.",
        "• Continuous Background Worker: Runs continuously in NestJS background (onModuleInit).",
        "• Semantic Token Deduplication: Overlap check (>50%) guarantees 100% unique problem solutions.",
        "• SOP Synthesis: Converts raw diagnostic work notes into structured ITIL Knowledge Articles."
    ]
    for b in k_bullets:
        pb = tf_k1.add_paragraph()
        pb.text = b
        pb.font.size = Pt(12)
        pb.font.color.rgb = TEXT_MUTED

    add_card(slide5, 6.8, 1.6, 5.7, 5.1)
    tb_k2 = slide5.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.7))
    tf_k2 = tb_k2.text_frame
    tf_k2.word_wrap = True
    
    p = tf_k2.paragraphs[0]
    p.text = "Synthesized SOP Article (KB0000001)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_AMBER

    kb_sample = """Title: Troubleshooting & SOP: SAP ERP Financials SSO Auth Failure (mainframe-host-01)
Category: Server - Kernel & OS Patch
Author: 🤖 Meta Llama 3.3 70B Instruct
Notes Analyzed: 64 Incidents

Resolution Steps:
1. Inspect configuration item telemetry status for mainframe-host-01.
2. Apply kernel remediation patch & clear SSO driver lock.
3. Validate metric health baseline and confirm ticket resolution."""

    pb = tf_k2.add_paragraph()
    pb.text = kb_sample
    pb.font.size = Pt(11)
    pb.font.color.rgb = TEXT_WHITE

    # ==================== SLIDE 6: MCP SERVER & WEB UI ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_bg(slide6)
    add_header(slide6, "MCP Server Protocol & Modern Next.js Frontend")

    add_card(slide6, 0.8, 1.6, 5.7, 5.1)
    tb_m1 = slide6.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.1), Inches(4.7))
    tf_m1 = tb_m1.text_frame
    tf_m1.word_wrap = True
    
    p = tf_m1.paragraphs[0]
    p.text = "Model Context Protocol (MCP)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_EMERALD

    m_bullets = [
        "• Transport: Stdio transport with Automatic Auth Flow.",
        "• Tool: list_unassigned_incidents (Query triage queue)",
        "• Tool: update_incident_group (Assign department & lead)",
        "• Tool: generate_knowledge_from_work_notes (Trigger SOP synthesis)",
        "• Tool: list_knowledge_articles (Retrieve published catalog)",
        "• Tool: get_knowledge_article (Inspect detailed SOP by ID)"
    ]
    for b in m_bullets:
        pb = tf_m1.add_paragraph()
        pb.text = b
        pb.font.size = Pt(12)
        pb.font.color.rgb = TEXT_MUTED

    add_card(slide6, 6.8, 1.6, 5.7, 5.1)
    tb_m2 = slide6.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.7))
    tf_m2 = tb_m2.text_frame
    tf_m2.word_wrap = True
    
    p = tf_m2.paragraphs[0]
    p.text = "Next.js Web UI (http://localhost:3000/knowledge)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_INDIGO

    ui_bullets = [
        "• Real-Time Progress Banner: Polling GET /api/v1/knowledge/status every 3s.",
        "• Zero State Loss: Navigating between pages (/incidents -> /knowledge) preserves progress.",
        "• Dynamic Controls: Trigger background worker across all 1,000 incidents.",
        "• Interactive SOP Modal: Deep-dive view of symptoms, root cause, and remediation steps."
    ]
    for b in ui_bullets:
        pb = tf_m2.add_paragraph()
        pb.text = b
        pb.font.size = Pt(12)
        pb.font.color.rgb = TEXT_WHITE

    # Output file
    output_path = os.path.join(os.getcwd(), "Agentic_AI_ITSM_Platform.pptx")
    prs.save(output_path)
    print(f"SUCCESS: Created PowerPoint presentation at {output_path}")

if __name__ == "__main__":
    build_pptx()
